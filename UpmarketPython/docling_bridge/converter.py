"""
Upmarket document converter — tiered pipeline.

Tier 1 (zero download): pdfium + post-processor for PDFs, markitdown for everything else
Tier 2 (172MB download): Enhanced pipeline — layout analysis for complex PDFs
Tier 3 (500MB download): Upmarket AI — Pro, scanned/research documents

No internal library names are exposed to callers.
"""

import os
import sys
import logging
import subprocess
from pathlib import Path
from docling_bridge.security import validate_file_path, validate_password, log_security_event

_ENHANCED_CONVERTERS = {}


def convert(file_path: str, options: dict | None = None) -> dict:
    """
    Convert a document to Markdown using the best available pipeline.

    Options:
        use_enhanced (bool): Use enhanced pipeline if downloaded. Default True.
        use_ai (bool): Use Upmarket AI (Pro). Default False.
        ocr (bool): Enable OCR. Default True.
        password (str): PDF password. Default None.

    Returns:
        {
            "success": bool,
            "markdown": str,
            "metadata": { "pages": int, "format": str, "title": str },
            "pipeline": str,  # "fast" | "enhanced" | "ai"
            "error": str | None,
            "needs_password": bool
        }
    """
    opts = options or {}

    # Security: validate path (traversal, size) and password (buffer overrun)
    try:
        file_path = validate_file_path(file_path)
        password_raw = opts.get("password", None)
        password_raw = validate_password(password_raw)
    except ValueError as e:
        log_security_event("INPUT_VALIDATION_FAILED", str(e))
        return _error(f"Upmarket couldn't open this file: {e}")

    path = Path(file_path)
    suffix = path.suffix.lower()
    password = password_raw
    use_enhanced = opts.get("use_enhanced", True)
    use_ai = opts.get("use_ai", False)

    # Check for password-protected PDF first
    if suffix == ".pdf" and not password:
        locked, err = _check_pdf_locked(str(path))
        if locked:
            return {**_error("This PDF is password-protected."), "needs_password": True}
        if err:
            return _error(err)

    # Formats Docling Enhanced supports (PDF + office documents only)
    # Audio, video, images always use fast path — Docling doesn't handle them
    ENHANCED_FORMATS = {'.pdf', '.docx', '.pptx', '.xlsx', '.html', '.htm',
                        '.md', '.asciidoc', '.epub', '.xml'}
    can_use_enhanced = suffix in ENHANCED_FORMATS
    can_use_ai = suffix in ENHANCED_FORMATS or suffix in {'.png', '.jpg', '.jpeg', '.tif', '.tiff', '.webp'}

    # Route to appropriate pipeline
    try:
        if use_ai and can_use_ai:
            if not _ai_available():
                return _error("Upmarket AI model is not downloaded or failed validation. Download it again from Settings > Models.")
            try:
                return _convert_ai(path, opts)
            except TimeoutError:
                raise
            except Exception as e:
                print(f"[Upmarket] AI pipeline failed for {suffix}: {e}", file=sys.stderr)
                return _error("Upmarket AI couldn't run on this Mac. Check model download and device compatibility.")

        if suffix == ".pdf" and use_enhanced and _enhanced_available():
            try:
                return _convert_enhanced(path, opts)
            except TimeoutError:
                raise
            except Exception as e:
                print(f"[Upmarket] Enhanced PDF fallback: {e}", file=sys.stderr)
                return _convert_fast_pdf(path, password)

        if suffix == ".pdf":
            return _convert_fast_pdf(path, password)

        # Non-PDF office formats: use enhanced if available and format is supported
        if use_enhanced and _enhanced_available() and can_use_enhanced:
            try:
                return _convert_enhanced(path, opts)
            except TimeoutError:
                raise
            except Exception as e:
                print(f"[Upmarket] Enhanced pipeline fallback for {suffix}: {e}", file=sys.stderr)
                return _convert_fast_other(path)

        # All other formats (audio, video, images, CSV, etc): always fast path
        return _convert_fast_other(path)

    except Exception as e:
        print(f"[Upmarket] Conversion error: {type(e).__name__}: {e}", file=sys.stderr)
        return _error("Upmarket couldn't convert this document. The file may be damaged, password-protected, or in an unsupported format.")


def check_pipelines() -> dict:
    """Returns which pipelines are available."""
    return {
        "fast": True,
        "enhanced": _enhanced_available(),
        "ai": _ai_available(),
    }


# MARK: - Pipeline implementations

def _convert_fast_pdf(path: Path, password: str | None) -> dict:
    """
    pdfium (Apache 2.0) + structured post-processing → clean Markdown.
    Uses block-level rect extraction with font-size heuristics for headings,
    paragraph merging, ligature fixing, and running header removal.
    """
    try:
        from docling_bridge.postprocessor import pdf_to_clean_markdown
        markdown, page_count = pdf_to_clean_markdown(str(path), password=password)
        return _success(markdown, page_count, path, pipeline="fast")
    except Exception as e:
        msg = str(e).lower()
        if "password" in msg or "encrypted" in msg:
            return {**_error("This PDF is password-protected."), "needs_password": True}
        print(f"[Upmarket] Post-processor error: {e}", file=sys.stderr)
        return _error("Upmarket couldn't convert this PDF.")


def _convert_fast_other(path: Path) -> dict:
    """
    Fast path for non-PDF formats.
    Routes by extension to the best available handler:
      - Images: handled in Swift with ImageIO before Python fallback
      - Audio/Video: handled in Swift with AVFoundation before Python fallback
      - Everything else: markitdown (MIT, Microsoft) — DOCX, PPTX, XLSX, HTML, CSV, EPUB
    """
    suffix = path.suffix.lower()
    speech_audio_suffixes = {'.mp3', '.m4a', '.wav'}

    # Images markitdown can't handle — use Pillow
    if suffix in ('.webp', '.tif', '.tiff', '.bmp', '.gif'):
        return _convert_image_pillow(path)

    # Native Swift handles first-choice audio transcription and media metadata.
    # Do not reject media here: MarkItDown may still provide useful Markdown for
    # supported audio files when the native route is unavailable.

    # Everything else: markitdown — with graceful partial-content fallback
    try:
        from markitdown import MarkItDown
        md_converter = MarkItDown()
        result = md_converter.convert(str(path))
        markdown = result.text_content or ""
        return _success(markdown, 0, path, pipeline="fast")
    except Exception as e:
        suffix_upper = path.suffix.upper().lstrip(".")
        err_msg = str(e)

        # Pillow fallback for any image format markitdown fails on
        if suffix in ('.png', '.jpg', '.jpeg'):
            return _convert_image_pillow(path)

        # DOCX with external images: mammoth fails on missing image references.
        # Try extracting text-only by stripping the problematic relationship.
        if suffix == '.docx' and ('KeyError' in err_msg or 'rId' in err_msg):
            return _convert_docx_text_only(path)

        # PPTX with unrecognized shapes: extract what we can, skip bad slides.
        if suffix == '.pptx' and 'NotImplementedError' in err_msg:
            return _convert_pptx_safe(path)

        if suffix in speech_audio_suffixes and "UnknownValueError" in err_msg:
            return _success(
                f"# Audio: {path.stem}\n\n_No recognizable speech was detected in this audio file._",
                1,
                path,
                pipeline="fast",
            )

        print(f"[Upmarket] markitdown error for {suffix_upper}: {e}", file=sys.stderr)
        return _error(f"Upmarket couldn't convert this {suffix_upper} file. Try downloading the Enhanced pipeline for better results.")


def _convert_docx_text_only(path: Path) -> dict:
    """Extract text from DOCX, skipping external images that cause mammoth to fail."""
    try:
        import zipfile
        import defusedxml.ElementTree as ET

        NS = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
        lines = []
        with zipfile.ZipFile(str(path)) as z:
            with z.open('word/document.xml') as f:
                tree = ET.parse(f)
                for para in tree.findall(f'.//{NS}p'):
                    text = ''.join(t.text or '' for t in para.findall(f'.//{NS}t'))
                    if text.strip():
                        lines.append(text.strip())

        return _success('\n\n'.join(lines), 0, path, pipeline="fast")
    except Exception as e:
        print(f"[Upmarket] DOCX text-only fallback failed: {e}", file=sys.stderr)
        return _error("Upmarket couldn't convert this Word document.")


def _convert_pptx_safe(path: Path) -> dict:
    """Extract text from PPTX, skipping slides with unrecognized shape types."""
    try:
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation(str(path))
        slides_md = []
        for i, slide in enumerate(prs.slides):
            lines = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'text') and shape.text.strip():
                        lines.append(shape.text.strip())
                except Exception:
                    pass  # skip unrecognized shapes
            if lines:
                slides_md.append(f"## Slide {i+1}\n\n" + '\n\n'.join(lines))

        return _success('\n\n'.join(slides_md), len(prs.slides), path, pipeline="fast")
    except Exception as e:
        print(f"[Upmarket] PPTX safe fallback failed: {e}", file=sys.stderr)
        return _error("Upmarket couldn't convert this PowerPoint file.")


def _convert_image_pillow(path: Path) -> dict:
    """Basic image fallback. Primary image metadata extraction is native Swift ImageIO."""
    try:
        from PIL import Image

        img = Image.open(str(path))
        width, height = img.size
        mode = img.mode
        fmt = img.format or path.suffix.lstrip('.').upper()

        lines = [
            f"# Image: {path.name}",
            f"",
            f"**Format:** {fmt}  ",
            f"**Dimensions:** {width} × {height} px  ",
            f"**Color mode:** {mode}  ",
        ]

        return _success("\n".join(lines), 1, path, pipeline="fast")
    except Exception as e:
        print(f"[Upmarket] Pillow error: {e}", file=sys.stderr)
        return _error(f"Upmarket couldn't read this image file.")


def _convert_media_native_only(path: Path) -> dict:
    """Deprecated compatibility shim. Media metadata is extracted with AVFoundation."""
    return _error(f"Upmarket couldn't read this media file.")


def _convert_enhanced(path: Path, opts: dict) -> dict:
    """Enhanced pipeline — handles all formats, complex layouts, tables."""
    from docling.document_converter import DocumentConverter, PdfFormatOption
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import PdfPipelineOptions

    pipeline_options = PdfPipelineOptions()
    pipeline_options.do_ocr = opts.get("ocr", True)
    pipeline_options.do_table_structure = True

    pdf_opts = PdfFormatOption(pipeline_options=pipeline_options)

    if opts.get("password"):
        try:
            pdf_opts = PdfFormatOption(
                pipeline_options=pipeline_options,
                backend_options={"password": opts["password"]}
            )
        except Exception:
            pass

    cache_key = None if opts.get("password") else (pipeline_options.do_ocr, pipeline_options.do_table_structure)
    converter = _ENHANCED_CONVERTERS.get(cache_key) if cache_key is not None else None
    if converter is None:
        format_options = {InputFormat.PDF: pdf_opts}
        converter = DocumentConverter(format_options=format_options)
        if cache_key is not None:
            _ENHANCED_CONVERTERS[cache_key] = converter

    result = converter.convert(str(path))
    markdown = result.document.export_to_markdown()
    page_count = result.document.num_pages() if callable(result.document.num_pages) else 0

    return _success(markdown, page_count, path, pipeline="enhanced")


_VLM_MAX_SIDE = 4096  # Granite MLX context limit; larger images return empty output


def _normalise_frame(img, max_side: int):
    """Convert a single PIL image frame to RGB and downsample if oversized."""
    from PIL import Image
    if img.mode in ("RGBA", "LA", "P"):
        img = img.convert("RGB")
    if max(img.size) > max_side:
        scale = max_side / max(img.size)
        img = img.resize((int(img.width * scale), int(img.height * scale)), Image.LANCZOS)
    return img


def _prepare_image_for_vlm(path: Path) -> Path:
    """Normalise an image for the VLM: flatten RGBA→RGB and downsample oversized images.

    Handles multi-page TIFFs by normalising every frame. Returns the original
    path unchanged when no normalisation is needed, or a temp path to a
    converted copy that the caller should treat as ephemeral.

    RGBA images cause the model to silently return empty Markdown; images wider
    or taller than _VLM_MAX_SIDE exhaust the VLM context window similarly.
    Multi-page TIFFs with RGBA frames timeout during MLX inference.
    """
    from PIL import Image
    import tempfile

    img = Image.open(path)
    suffix = path.suffix.lower()
    is_tiff = suffix in (".tif", ".tiff")

    # Collect all frames for multi-page TIFFs, single frame otherwise.
    frames = []
    try:
        while True:
            frames.append(img.copy())
            img.seek(img.tell() + 1)
    except EOFError:
        pass

    needs_normalise = any(
        f.mode in ("RGBA", "LA", "P") or max(f.size) > _VLM_MAX_SIDE
        for f in frames
    )

    if not needs_normalise:
        return path

    normalised = [_normalise_frame(f, _VLM_MAX_SIDE) for f in frames]
    w0, h0 = frames[0].size

    fmt = "TIFF" if is_tiff else ("JPEG" if suffix in (".jpg", ".jpeg") else "PNG")
    tmp = tempfile.NamedTemporaryFile(suffix=f".{suffix or '.png'}", delete=False)
    if len(normalised) > 1:
        normalised[0].save(tmp.name, format=fmt, save_all=True, append_images=normalised[1:])
    else:
        normalised[0].save(tmp.name, format=fmt)

    print(
        f"[Upmarket] VLM image prepared: {path.name} {len(frames)}p {w0}x{h0}"
        f" → {normalised[0].width}x{normalised[0].height} RGB",
        file=sys.stderr,
    )
    return Path(tmp.name)


def _convert_ai(path: Path, opts: dict) -> dict:
    """Upmarket AI — Pro tier, Granite Docling MLX path for image/scanned documents."""
    manager = _model_manager()
    if not manager.supports_upmarket_ai_hardware():
        return _error("Upmarket AI requires Apple Silicon with Metal support.")
    runtime_error = _upmarket_ai_runtime_unavailable_reason()
    if runtime_error:
        return _error(runtime_error)

    suffix = path.suffix.lower()
    if suffix not in (".pdf", ".png", ".jpg", ".jpeg", ".tif", ".tiff", ".webp"):
        return _convert_enhanced(path, opts)

    model_path = _ai_model_path()
    artifacts_path = model_path.parent
    _quiet_known_granite_warnings()

    # Import the VLM pipeline only after model validation. In headless/test
    # environments MLX may not see a Metal device, and that must not break
    # non-AI conversion.
    from docling.datamodel import vlm_model_specs
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import VlmPipelineOptions
    from docling.document_converter import DocumentConverter, ImageFormatOption, PdfFormatOption
    from docling.pipeline.vlm_pipeline import VlmPipeline

    pipeline_options = VlmPipelineOptions(
        artifacts_path=artifacts_path,
        vlm_options=vlm_model_specs.GRANITEDOCLING_MLX,
    )
    if suffix == ".pdf":
        format_options = {
            InputFormat.PDF: PdfFormatOption(
                pipeline_cls=VlmPipeline,
                pipeline_options=pipeline_options,
            )
        }
    else:
        format_options = {
            InputFormat.IMAGE: ImageFormatOption(
                pipeline_cls=VlmPipeline,
                pipeline_options=pipeline_options,
            )
        }

    converter = DocumentConverter(format_options=format_options)
    convert_path = _prepare_image_for_vlm(path) if suffix != ".pdf" else path
    result = converter.convert(convert_path)
    markdown = result.document.export_to_markdown()
    if not markdown.strip():
        raise RuntimeError("AI model returned empty Markdown")
    page_count = result.document.num_pages() if callable(result.document.num_pages) else 0
    return _success(markdown, page_count, path, pipeline="ai")


class _GraniteNoiseFilter(logging.Filter):
    def filter(self, record: logging.LogRecord) -> bool:
        message = record.getMessage()
        return not (
            "Model config: pad_token_id must be `None` or an integer within the vocabulary" in message
            and "got 128002" in message
        )


_GRANITE_WARNINGS_QUIETED = False
_AI_RUNTIME_PRECHECK: str | None = None

_AI_RUNTIME_PROBE_SCRIPT = """
import sys

try:
    import mlx.core as mx
    if not mx.device_info():
        print("empty graphics processor response", file=sys.stderr)
        raise SystemExit(2)
except Exception as exc:
    print(str(exc), file=sys.stderr)
    raise SystemExit(1)

raise SystemExit(0)
"""


def _quiet_known_granite_warnings() -> None:
    global _GRANITE_WARNINGS_QUIETED
    if _GRANITE_WARNINGS_QUIETED:
        return

    logging.getLogger("transformers.configuration_utils").addFilter(_GraniteNoiseFilter())

    try:
        import mlx.core as mx

        if hasattr(mx, "device_info") and hasattr(mx, "metal"):
            mx.metal.device_info = mx.device_info
    except Exception:
        pass

    _GRANITE_WARNINGS_QUIETED = True


def _upmarket_ai_runtime_unavailable_reason() -> str | None:
    """Return a user-safe reason when the current session cannot run Upmarket AI."""
    global _AI_RUNTIME_PRECHECK
    if _AI_RUNTIME_PRECHECK is not None:
        return None if _AI_RUNTIME_PRECHECK == "" else _AI_RUNTIME_PRECHECK

    message = "Upmarket AI cannot access this Mac's graphics processor from the current session. Quit and reopen Upmarket, then try again."
    env = os.environ.copy()
    env["HF_HUB_OFFLINE"] = "1"
    env["TRANSFORMERS_OFFLINE"] = "1"

    try:
        completed = subprocess.run(
            [sys.executable, "-c", _AI_RUNTIME_PROBE_SCRIPT],
            capture_output=True,
            text=True,
            timeout=15,
            env=env,
        )
        if completed.returncode == 0:
            _AI_RUNTIME_PRECHECK = ""
        else:
            detail = (completed.stderr or completed.stdout or f"exit {completed.returncode}").strip()
            print(f"[Upmarket] AI runtime preflight failed: {detail[-300:]}", file=sys.stderr)
            _AI_RUNTIME_PRECHECK = message
    except (OSError, subprocess.TimeoutExpired) as exc:
        print(f"[Upmarket] AI runtime preflight failed: {exc}", file=sys.stderr)
        _AI_RUNTIME_PRECHECK = message

    return None if _AI_RUNTIME_PRECHECK == "" else _AI_RUNTIME_PRECHECK


def _ai_model_path() -> Path:
    manager = _model_manager()

    model_path = manager.model_directory("upmarket_ai")
    ok, error = manager.validate_model_dir("upmarket_ai", model_path)
    if not ok:
        raise RuntimeError(error or "Upmarket AI model is not available")
    return model_path.resolve()


# MARK: - Availability checks

def _enhanced_available() -> bool:
    try:
        return _model_manager().model_available("layout")
    except Exception:
        return False


def _ai_available() -> bool:
    try:
        return _model_manager().model_available("upmarket_ai")
    except Exception:
        return False


def _model_manager():
    try:
        from models import model_manager
    except ImportError:
        from upmarket_models import model_manager
    return model_manager


# MARK: - Helpers

def _check_pdf_locked(file_path: str) -> tuple[bool, str | None]:
    try:
        import pypdfium2 as pdfium
        doc = pdfium.PdfDocument(file_path)
        doc.close()
        return False, None
    except Exception as e:
        msg = str(e).lower()
        if "password" in msg or "encrypted" in msg:
            return True, None
        return False, None


def _success(markdown: str, pages: int, path: Path, pipeline: str) -> dict:
    return {
        "success": True,
        "markdown": markdown,
        "metadata": {
            "pages": pages,
            "format": path.suffix.lstrip(".").upper(),
            "title": path.stem,
        },
        "pipeline": pipeline,
        "error": None,
        "needs_password": False,
    }


def _error(message: str) -> dict:
    return {
        "success": False,
        "markdown": "",
        "metadata": {},
        "pipeline": "none",
        "error": message,
        "needs_password": False,
    }
